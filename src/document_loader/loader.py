"""
Document Loader Module
Supports loading PDF, PPT, DOCX, and Image files
"""

import re
from abc import ABC, abstractmethod     #抽象基类
from pathlib import Path                #处理文件路径
from typing import List, Optional       #类型提示
from dataclasses import dataclass       #数据类，装饰器使用

#Document类
#读取后的文档格式
@dataclass
class Document:
    """Represents a loaded document with metadata"""
    content: str           #文档里面的文字内容
    file_path: str         #文件路径
    file_type: str         #文件类型（pdf/pot/docx/image）
    metadata: dict         #额外信息（页数，幻灯片数..）

#BaseLoader 抽象类（模板）
class BaseLoader(ABC):
    """Base class for document loaders"""

    #读取文件
    @abstractmethod
    def load(self, file_path: str) -> Document:
        """Load a document from the given file path"""
        pass
    #判断是否支持该类型
    @abstractmethod
    def supports(self, file_extension: str) -> bool:
        """Check if this loader supports the given file extension"""
        pass

#PDFLoader读取PDF
class PDFLoader(BaseLoader):
    """Loader for PDF files"""

    #接收一个文件路径，返回一个Document对象
    def load(self, file_path: str) -> Document:
        try:
            import pdfplumber
        except ImportError:
            try:
                from pypdf import PdfReader
                pdfplumber_available = False
            except ImportError:
                raise ImportError("Please install either pdfplumber or pypdf: pip install pdfplumber")

        content = []
        metadata = {}

        # Try pdfplumber first (better for text extraction and tables)
        if 'pdfplumber' in globals() or 'pdfplumber' in locals() or True:  # Always try to use it if available
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    pages_text = []
                    tables_data = []

                    for i, page in enumerate(pdf.pages):
                        # Extract text from page
                        page_text = page.extract_text()
                        if page_text:
                            pages_text.append(page_text)

                        # Extract tables from page
                        tables = page.extract_tables()
                        for table in tables:
                            if table:
                                # Convert table to text format
                                table_text = "\n".join(["\t".join([cell or "" for cell in row]) for row in table])
                                tables_data.append(table_text)

                    content = pages_text
                    metadata["num_pages"] = len(pdf.pages)
                    metadata["tables_extracted"] = len(tables_data)
                    if tables_data:
                        content.extend(tables_data)
            except ImportError:
                # Fallback to pypdf
                from pypdf import PdfReader
                reader = PdfReader(file_path)
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        content.append(text)

                metadata["num_pages"] = len(reader.pages)

        return Document(
            content="\n".join(content),   #将每页的文本片段连接成一段长文本，使用\n隔离
            file_path=file_path,
            file_type="pdf",
            metadata=metadata
        )

    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() == ".pdf"

#读取ppt
class PPTLoader(BaseLoader):
    “””Loader for PPT/PPTX files”””   #加载ppt文件

    def load(self, file_path: str) -> Document:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(“Please install python-pptx: pip install python-pptx”)

        #创建ppt加载起对象
        prs = Presentation(file_path)
        content = []
        tables_data = []

        #prs.slides 幻灯片数  i索引 slide幻灯片对象
        for i, slide in enumerate(prs.slides):
            slide_text = []
            #遍历当前幻灯片的所有元素
            for shape in slide.shapes:
                #检查幻灯片是否有”text”元素，并且text元素是否不为空
                if hasattr(shape, “text”) and shape.text:
                    slide_text.append(shape.text)

                # Extract table data if it exists
                if hasattr(shape, “has_table”) and shape.has_table:
                    table = shape.table
                    table_content = []
                    for row in table.rows:
                        row_content = []
                        for cell in row.cells:
                            row_content.append(cell.text)
                        table_content.append(“\t”.join(row_content))
                    tables_data.append(“\n”.join(table_content))

            if slide_text:
                #提取当前幻灯片的text的内容，[Slide 页数{i+1}]\n加上一个幻灯片的
                content.append(f”[Slide {i+1}]\n” + “\n”.join(slide_text))

        if tables_data:
            content.extend(tables_data)

        return Document(
            content=”\n”.join(content),
            file_path=file_path,
            file_type=”ppt”,
            metadata={“num_slides”: len(prs.slides), “tables_extracted”: len(tables_data)}
        )

    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in [“.ppt”, “.pptx”]

#读取docx
class DOCXLoader(BaseLoader):
    """Loader for DOCX files"""

    def load(self, file_path: str) -> Document:
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise ImportError("Please install python-docx: pip install python-docx")

        doc = DocxDocument(file_path)
        #将word文档的非空段落提取出来，存成一个干净的文本列表
        paragraphs = [para.text for para in doc.paragraphs if para.text]

        # Extract tables if they exist
        tables_data = []
        for table in doc.tables:
            table_content = []
            for row in table.rows:
                row_content = []
                for cell in row.cells:
                    row_content.append(cell.text)
                table_content.append("\t".join(row_content))
            if table_content:
                tables_data.append("\n".join(table_content))

        content = paragraphs[:]
        if tables_data:
            content.extend(tables_data)

        return Document(
            content="\n".join(content),
            file_path=file_path,
            file_type="docx",
            metadata={"num_paragraphs": len(paragraphs), "tables_extracted": len(tables_data)}
        )

    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() == ".docx"

#图片加载器 核心能力是OCR（光学字符识别）
class ImageLoader(BaseLoader):
    """Loader for Image files with OCR support"""
    #参数lang  chi_simt：同时识别中文和英文
    def __init__(self, lang: str = "chi_sim+eng"):
        self.lang = lang

    def load(self, file_path: str) -> Document:
        try:
            import pytesseract            #将图片转成文字
            from PIL import Image         #图像处理库，打开图片文件，将图片数据交给Tesseract
        except ImportError:
            raise ImportError("Please install pytesseract and Pillow: pip install pytesseract Pillow")

        # Try using Unstructured for more robust image processing if available
        try:
            from unstructured.partition.auto import partition
            from unstructured.documents.elements import Text
            element = partition(filename=file_path)
            unstructured_text = "\n".join([str(el) for el in element if isinstance(el, Text)])

            # If unstructured works well, use it; otherwise fallback to Tesseract
            if unstructured_text and len(unstructured_text.strip()) > 10:
                content = unstructured_text
            else:
                # Fallback to Tesseract
                image = Image.open(file_path)
                content = pytesseract.image_to_string(image, lang=self.lang)        #image传入图片，  设置中英文  转换成字符串存入text
        except ImportError:
            # Unstructured not available, use Tesseract
            from PIL import Image
            image = Image.open(file_path)
            content = pytesseract.image_to_string(image, lang=self.lang)

        return Document(
            content=content,
            file_path=file_path,
            file_type="image",
            metadata={"image_size": image.size if 'image' in locals() else (0, 0), "mode": image.mode if 'image' in locals() else "unknown"}
        )

    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".gif"]


class CSVLoader(BaseLoader):
    """Loader for CSV files with support for tabular data extraction"""

    def load(self, file_path: str) -> Document:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("Please install pandas: pip install pandas")

        df = pd.read_csv(file_path)

        # Convert DataFrame to text representation
        csv_text = df.to_string(index=False)

        # Also include a structured summary
        summary_lines = [
            f"CSV File Summary:",
            f"Rows: {len(df)}",
            f"Columns: {len(df.columns)}",
            f"Column Names: {', '.join(df.columns.tolist())}",
            "",
            "Data:",
            csv_text
        ]

        return Document(
            content="\n".join(summary_lines),
            file_path=file_path,
            file_type="csv",
            metadata={
                "num_rows": len(df),
                "num_columns": len(df.columns),
                "column_names": df.columns.tolist()
            }
        )

    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() == ".csv"


class ExcelLoader(BaseLoader):
    """Loader for Excel files with multi-sheet support"""

    def load(self, file_path: str) -> Document:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("Please install pandas: pip install pandas")

        content_parts = []
        metadata = {"sheets": []}

        # Load all sheets
        excel_file = pd.ExcelFile(file_path)

        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)

            # Convert sheet to text
            sheet_text = f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n"
            content_parts.append(sheet_text)

            metadata["sheets"].append({
                "name": sheet_name,
                "rows": len(df),
                "columns": len(df.columns)
            })

        content = "\n".join(content_parts)

        return Document(
            content=content,
            file_path=file_path,
            file_type="excel",
            metadata=metadata
        )

    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in [".xlsx", ".xls"]


class DocumentCleaner:
    """Clean and preprocess document content"""

    @staticmethod
    def clean_text(text: str) -> str:
        """
        Clean and normalize text content
        - Remove headers/footers
        - Clean up whitespace
        - Remove common document artifacts
        """
        if not text:
            return text

        # Remove multiple consecutive blank lines
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)

        # Remove page numbers that appear at beginning or end of lines
        text = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)
        text = re.sub(r'^\s*[-*_]{3,}\s*$', '', text, flags=re.MULTILINE)

        # Remove header/footer artifacts (patterns like document titles repeated)
        lines = text.split('\n')
        if len(lines) > 10:  # Only if document is substantial
            # Remove lines that appear frequently at the beginning/end
            first_few_lines = lines[:5]
            last_few_lines = lines[-5:]

            # Count occurrences of lines to identify header/footer patterns
            from collections import Counter
            all_lines_counter = Counter(lines)

            # Remove lines that appear too frequently (likely headers/footers)
            cleaned_lines = []
            for line in lines:
                if all_lines_counter[line] <= 2 or len(line.strip()) < 50:  # Keep short lines
                    cleaned_lines.append(line)

            text = '\n'.join(cleaned_lines)

        # Remove extra whitespace but preserve sentence structure
        text = re.sub(r'[ \t]+', ' ', text)  # Multiple spaces/tabs to single space
        text = re.sub(r' +\n', '\n', text)  # Space before newline
        text = re.sub(r'\n +', '\n', text)  # Space after newline

        return text.strip()


class DocumentLoader:
    """
    Unified document loader that automatically selects the appropriate loader
    based on file extension
    """
    #在构造构造函数初始化这些工具类    ->设计模式  策略模式类似于工厂模式 好处：一次配置，终身受用
    def __init__(self):
        self.loaders: List[BaseLoader] = [
            PDFLoader(),
            PPTLoader(),
            DOCXLoader(),
            CSVLoader(),
            ExcelLoader(),
            ImageLoader()
        ]
    #传入文件路径，判断是用哪个文件解析器
    def load(self, file_path: str) -> Document:
        """Load a document from the given file path"""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        #获取文件后缀名
        extension = path.suffix
        for loader in self.loaders:
            #分发，看它支持哪个文件格式
            if loader.supports(extension):
                doc = loader.load(file_path)
                # Clean the content after loading
                cleaned_content = DocumentCleaner.clean_text(doc.content)
                return Document(
                    content=cleaned_content,
                    file_path=doc.file_path,
                    file_type=doc.file_type,
                    metadata=doc.metadata
                )

        raise ValueError(f"Unsupported file type: {extension}")
        #传入目录路径
    def load_directory(self, directory_path: str, file_types: Optional[List[str]] = None) -> List[Document]:
        """Load all documents from a directory"""
        path = Path(directory_path)
        if not path.is_dir():
            raise NotADirectoryError(f"Not a directory: {directory_path}")

        documents = []
        #path.iterdir  这是一个生成器，他会一个一个列出文件夹的所有内容（包括子文件夹和文件）
        for file_path in path.iterdir():
            if file_types is None or file_path.suffix in file_types:   #如果没在指定类型全部接收，如果指定了就选择指定的文件
                try:
                    #传给文档解析器
                    doc = self.load(str(file_path))
                    #将文档解析起解析的内容 追加document
                    documents.append(doc)
                except (ValueError, ImportError) as e:
                    print(f"Skipping {file_path}: {e}")

        return documents


if __name__ == "__main__":
    # Example usage
    loader = DocumentLoader()
    doc = loader.load("example.pdf")
    print(f"Loaded {doc.file_type} with {len(doc.content)} characters")
